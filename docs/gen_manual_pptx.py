#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# Manual de usuario v2 del Bot WhatsApp Grupo Ardisa — PowerPoint PREMIUM (rediseño pedido Deicy 12-ago).
# Diseño: papel claro con número fantasma por sección, kicker dorado, subrayado esmeralda, tarjetas con
# insignias circulares, mock de chat estilo WhatsApp dibujado con formas, tablas altas con primera columna
# en negrita. Portada y cierre en azul marino con burbuja decorativa.
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x0E, 0x2A, 0x3B)
NAVY2  = RGBColor(0x14, 0x32, 0x46)   # burbuja decorativa sobre navy
NAVY3  = RGBColor(0x1B, 0x3E, 0x54)
EMER   = RGBColor(0x0E, 0x8A, 0x4E)
EMERD  = RGBColor(0x0A, 0x6B, 0x3C)
GOLD   = RGBColor(0xC9, 0x9B, 0x3F)
GOLDD  = RGBColor(0x9A, 0x6E, 0x10)
PAPER  = RGBColor(0xF7, 0xF8, 0xF4)
GHOST  = RGBColor(0xEC, 0xF1, 0xEA)   # número fantasma (muy tenue)
INK    = RGBColor(0x18, 0x24, 0x20)
SOFT   = RGBColor(0x5C, 0x6E, 0x63)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xDF, 0xE7, 0xDF)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
CHATBG = RGBColor(0xEA, 0xF0, 0xE6)
BUBBOT = RGBColor(0xD8, 0xF2, 0xDF)
BUBPAN = RGBColor(0xCD, 0xEB, 0xF7)
WARNBG = RGBColor(0xFD, 0xF6, 0xE6)
WARNBD = RGBColor(0xE3, 0xC8, 0x82)
WARNTX = RGBColor(0x5B, 0x43, 0x00)
ALTROW = RGBColor(0xF0, 0xF5, 0xF0)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height
F = 'Segoe UI'

def fondo(s, color=PAPER):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color

def rect(s, x, y, w, h, fill=None, line=None, radio=0.08, forma=MSO_SHAPE.ROUNDED_RECTANGLE, lw=1.0):
    sh = s.shapes.add_shape(forma, x, y, w, h)
    try: sh.adjustments[0] = radio
    except Exception: pass
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh

def tx(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=6):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(paras if isinstance(paras, list) else [paras]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, st in (para if isinstance(para, list) else [(para, {})]):
            r = p.add_run(); r.text = t
            r.font.size = Pt(st.get('s', 15)); r.font.bold = st.get('b', False)
            r.font.color.rgb = st.get('c', INK); r.font.name = st.get('f', F)
            if st.get('i'): r.font.italic = True
    return tb

def cab(s, num, titulo, sub=None):
    """Encabezado premium: número fantasma + kicker dorado + título con subrayado esmeralda."""
    fondo(s)
    tx(s, Inches(11.55), Inches(0.16), Inches(1.55), Inches(1.05),
       [[(num, {'s': 60, 'b': True, 'c': GHOST})]], align=PP_ALIGN.RIGHT)
    tx(s, Inches(0.65), Inches(0.42), Inches(9.5), Inches(0.32),
       [[("M A N U A L   D E   U S U A R I O   ·   S E C C I Ó N   " + num, {'s': 10.5, 'b': True, 'c': GOLDD})]])
    tx(s, Inches(0.62), Inches(0.68), Inches(10.2), Inches(0.75),
       [[(titulo, {'s': 30, 'b': True, 'c': NAVY})]])
    rect(s, Inches(0.68), Inches(1.42), Inches(1.5), Inches(0.055), fill=EMER, radio=0.5)
    if sub:
        tx(s, Inches(2.4), Inches(1.28), Inches(10.2), Inches(0.34), [[(sub, {'s': 12.5, 'c': SOFT, 'i': True})]])

def pie(s, n):
    rect(s, Inches(0.65), Inches(7.06), Inches(12.03), Pt(1).emu and Inches(0.012), fill=LINE, radio=0.5)
    tx(s, Inches(0.65), Inches(7.12), Inches(12.03), Inches(0.3),
       [[("Bot WhatsApp Grupo Ardisa · Manual de usuario v2 · agosto 2026", {'s': 9.5, 'c': SOFT}),
         ("        %02d" % n, {'s': 9.5, 'c': GOLDD, 'b': True})]])

def tabla(s, x, y, w, filas, anchos, alto=0.62, size=12.5):
    gt = s.shapes.add_table(len(filas), len(filas[0]), x, y, w, Inches(alto * len(filas))).table
    for j, a in enumerate(anchos): gt.columns[j].width = Inches(a)
    for i, fila in enumerate(filas):
        for j, val in enumerate(fila):
            c = gt.cell(i, j); c.text = val
            c.margin_left = Inches(0.12); c.margin_right = Inches(0.1)
            c.margin_top = Inches(0.05); c.margin_bottom = Inches(0.05)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = F; r.font.size = Pt(size)
                    if i == 0: r.font.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(size - 1)
                    elif j == 0: r.font.bold = True; r.font.color.rgb = NAVY
                    else: r.font.color.rgb = INK
            c.fill.solid()
            c.fill.fore_color.rgb = NAVY if i == 0 else (WHITE if i % 2 else ALTROW)
    return gt

def nota(s, x, y, w, txt_, emoji="⚠️", h=0.8):
    cj = rect(s, x, y, w, Inches(h), fill=WARNBG, line=WARNBD, radio=0.14)
    tf = cj.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.16)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = emoji + "   " + txt_
    r.font.size = Pt(12.5); r.font.color.rgb = WARNTX; r.font.name = F
    return cj

def badge(s, x, y, num):
    b = rect(s, x, y, Inches(0.42), Inches(0.42), fill=EMER, forma=MSO_SHAPE.OVAL)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = F

def burbuja(s, x, y, w, texto_, lado='cli', quien=None, size=10.5):
    """Burbuja estilo WhatsApp; x/y/w EN PULGADAS (float). Devuelve la y siguiente, tambien en pulgadas."""
    fill = WHITE if lado == 'cli' else (BUBBOT if lado == 'bot' else BUBPAN)
    lineas = max(1, int(len(texto_) / 34) + 1) + (1 if quien else 0)
    h = 0.22 + 0.185 * lineas
    cj = rect(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=fill, line=LINE, radio=0.16)
    tf = cj.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.13); tf.margin_right = Inches(0.11); tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    if quien:
        r = p.add_run(); r.text = quien + "\n"
        r.font.size = Pt(8.5); r.font.bold = True; r.font.name = F
        r.font.color.rgb = EMERD if lado == 'bot' else RGBColor(0x0B, 0x5D, 0x74)
    r = p.add_run(); r.text = texto_
    r.font.size = Pt(size); r.font.color.rgb = INK; r.font.name = F
    return y + h + 0.09

# ═════════ 1 · PORTADA ═════════
s = prs.slides.add_slide(BLANK); fondo(s, NAVY)
rect(s, Inches(8.6), Inches(1.1), Inches(4.9), Inches(3.4), fill=NAVY2, radio=0.18)
rect(s, Inches(7.9), Inches(3.6), Inches(3.4), Inches(2.2), fill=NAVY3, radio=0.22)
rect(s, Inches(11.6), Inches(5.55), Inches(0.5), Inches(0.5), fill=EMER, forma=MSO_SHAPE.OVAL)
tx(s, Inches(8.6), Inches(1.7), Inches(4.9), Inches(2.2), [[("💬", {'s': 96, 'c': WHITE})]], align=PP_ALIGN.CENTER)
tx(s, Inches(8.0), Inches(4.05), Inches(3.2), Inches(1.4),
   [[("«Hola, necesito 30\nbultos de cemento»", {'s': 13, 'c': RGBColor(0xBE, 0xD4, 0xCE), 'i': True})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
tx(s, Inches(0.75), Inches(1.35), Inches(7.6), Inches(0.4),
   [[("G R U P O   A R D I S A   ·   L Í N E A   3 1 6", {'s': 13, 'b': True, 'c': GOLD})]])
tx(s, Inches(0.72), Inches(1.85), Inches(7.9), Inches(2.2),
   [[("Manual de usuario", {'s': 47, 'b': True, 'c': WHITE})],
    [("del bot y su panel", {'s': 47, 'b': True, 'c': GOLD})]], sp=0)
rect(s, Inches(0.78), Inches(4.05), Inches(2.0), Inches(0.05), fill=EMER, radio=0.5)
tx(s, Inches(0.75), Inches(4.35), Inches(6.7), Inches(1.7),
   [[("Versión 2 · agosto de 2026. El bot con inteligencia artificial (Fase 1 + 2), el panel de conversaciones, el seguimiento de solicitudes, las alertas del vigilante y qué hacer cuando algo falla.", {'s': 15, 'c': RGBColor(0xC9, 0xDA, 0xD5)})]])
tx(s, Inches(0.75), Inches(6.55), Inches(11), Inches(0.4),
   [[("Para el equipo de coordinación · Reemplaza el manual de la Fase 1", {'s': 11, 'c': RGBColor(0x8F, 0xA8, 0xA2)})]])

# ═════════ 2 · CONTENIDO ═════════
s = prs.slides.add_slide(BLANK); cab(s, "☰", "Contenido")
SECS = [("01", "El sistema en un mapa"), ("02", "Panel de conversaciones"), ("03", "Chat híbrido (instalado, apagado)"),
        ("04", "Seguimiento de solicitudes"), ("05", "Tu panel por WhatsApp"), ("06", "Cómo reportan los asesores"),
        ("07", "Alertas del vigilante"), ("08", "Correos automáticos"), ("09", "Fase 2: la IA y SAP"), ("10", "Si algo falla")]
for idx, (n_, t_) in enumerate(SECS):
    col, fila = idx % 2, idx // 2
    x = Inches(0.75 + col * 6.2); y = Inches(1.95 + fila * 0.95)
    rect(s, x, y, Inches(5.9), Inches(0.78), fill=CARD, line=LINE, radio=0.14)
    tx(s, x + Inches(0.22), y + Inches(0.13), Inches(0.9), Inches(0.5), [[(n_, {'s': 20, 'b': True, 'c': GOLDD})]])
    tx(s, x + Inches(1.05), y + Inches(0.16), Inches(4.7), Inches(0.5), [[(t_, {'s': 16, 'b': True, 'c': NAVY})]])
pie(s, 2)

# ═════════ 3 · MAPA ═════════
s = prs.slides.add_slide(BLANK); cab(s, "01", "El sistema en un mapa", "cinco piezas que trabajan juntas")
tabla(s, Inches(0.65), Inches(1.8), Inches(12.05), [
    ["Pieza", "Qué hace", "Dónde la ves"],
    ["El bot", "Atiende el 316 las 24 horas: autorización de datos, entiende el pedido con IA, rutea al asesor experto y registra todo.", "El cliente lo vive en WhatsApp"],
    ["Panel de conversaciones", "Cada chat como en WhatsApp: texto, fotos, audios y documentos.", "n8n.ardisa.com/monitor/"],
    ["Seguimiento", "La tabla de solicitudes con estados, valores y filtros; exporta a Excel.", "…/monitor/seguimiento.php"],
    ["Vigilante", "Revisa el bot cada hora: clientes perdidos, crones caídos, colas atascadas.", "Correo a Deicy + panel de WhatsApp"],
    ["Reportes automáticos", "Excel semanal por marca a coordinación; chequeo de duplicados diario.", "Correo, lunes 7:00 a.m."],
], [2.35, 6.6, 3.1], alto=0.74, size=12.5)
nota(s, Inches(0.65), Inches(6.15), Inches(12.05),
     "El panel solo abre desde la red de la oficina (MikroTik) o la VPN. Si no abre desde tu casa, esa es la razón — no es un daño.", "🔒", 0.7)
pie(s, 3)

# ═════════ 4 · PANEL ═════════
s = prs.slides.add_slide(BLANK); cab(s, "02", "Panel de conversaciones", "n8n.ardisa.com/monitor/")
items = [
    ("👥", "Clientes y asesores, separados", "Los chats donde los asesores reportan van en su propia pestaña, sin mezclarse con los clientes."),
    ("📆", "Hoy · Ayer · Todos", "El contador dice cuántos chats hay y cuántos están 🟢 en curso (actividad en los últimos 30 minutos)."),
    ("🗑️", "Ocultar sin borrar", "Saca de la lista el spam o los proveedores; los mensajes quedan intactos y se restauran desde «Ver ocultos»."),
    ("📎", "Todo se ve", "Fotos, audios, videos y documentos del cliente, dentro del chat, como en WhatsApp."),
    ("🔄", "Se actualiza sola", "La página se refresca cada 15 segundos, sin que tengas que hacer nada."),
]
y = 1.85
for e, t_, d_ in items:
    rect(s, Inches(0.75), Inches(y), Inches(11.8), Inches(0.88), fill=CARD, line=LINE, radio=0.12)
    tx(s, Inches(1.0), Inches(y + 0.17), Inches(0.6), Inches(0.5), [[(e, {'s': 21})]])
    tx(s, Inches(1.72), Inches(y + 0.09), Inches(3.4), Inches(0.66), [[(t_, {'s': 14.5, 'b': True, 'c': NAVY})]], anchor=MSO_ANCHOR.MIDDLE)
    tx(s, Inches(5.25), Inches(y + 0.09), Inches(7.1), Inches(0.66), [[(d_, {'s': 12.5, 'c': SOFT})]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.0
pie(s, 4)

# ═════════ 5 · CHAT HÍBRIDO ═════════
s = prs.slides.add_slide(BLANK); cab(s, "03", "Chat híbrido: responder desde el panel")
rect(s, Inches(9.05), Inches(0.52), Inches(3.62), Inches(0.42), fill=WARNBG, line=WARNBD, radio=0.5)
tx(s, Inches(9.14), Inches(0.575), Inches(3.5), Inches(0.34), [[("INSTALADO · APAGADO POR AHORA", {'s': 10.5, 'b': True, 'c': WARNTX})]])
# mock de chat
rect(s, Inches(0.72), Inches(1.75), Inches(4.55), Inches(4.9), fill=CHATBG, line=LINE, radio=0.07)
yb = 2.0
yb = burbuja(s, 0.95, yb, 3.3, "Buenas, ¿me confirmas el despacho de mi pedido?", 'cli', "Cliente")
yb = burbuja(s, 1.75, yb, 3.4, "Tu solicitud ya está en gestión con nuestra asesora. 🤝", 'bot', "🤖 Bot")
yb = burbuja(s, 1.75, yb, 3.4, "¡Hola! Te confirmo: tu pedido sale hoy en la tarde. 🙌", 'pan', "👩‍💼 Ardisa (panel)")
rect(s, Inches(0.95), Inches(6.0), Inches(3.5), Inches(0.44), fill=WHITE, line=LINE, radio=0.5)
tx(s, Inches(1.12), Inches(6.06), Inches(2.9), Inches(0.32), [[("Escribe una respuesta…", {'s': 10.5, 'c': SOFT, 'i': True})]])
rect(s, Inches(4.52), Inches(6.0), Inches(0.44), Inches(0.44), fill=EMER, forma=MSO_SHAPE.OVAL)
tx(s, Inches(4.55), Inches(6.045), Inches(0.4), Inches(0.36), [[("➤", {'s': 13, 'c': WHITE, 'b': True})]], align=PP_ALIGN.CENTER)
# pasos a la derecha
pasos = [("1", "Tomas la conversación", "Botón «👩‍💼 Atender yo» (o simplemente respondes). El bot se calla para ESE cliente; los demás siguen normal."),
         ("2", "Escribes por el número del bot", "Enter envía. Tu burbuja sale azul, «Ardisa (panel)», y todo queda en el historial."),
         ("3", "Devuelves al bot", "Botón «🤖 Devolver al bot», o el bot retoma solo a los 30 minutos de tu última respuesta.")]
y = 1.85
for n_, t_, d_ in pasos:
    rect(s, Inches(5.7), Inches(y), Inches(6.95), Inches(1.06), fill=CARD, line=LINE, radio=0.12)
    badge(s, Inches(5.95), Inches(y + 0.31), n_)
    tx(s, Inches(6.6), Inches(y + 0.1), Inches(5.9), Inches(0.4), [[(t_, {'s': 14.5, 'b': True, 'c': NAVY})]])
    tx(s, Inches(6.6), Inches(y + 0.46), Inches(5.9), Inches(0.55), [[(d_, {'s': 11.5, 'c': SOFT})]])
    y += 1.22
nota(s, Inches(5.7), Inches(5.6), Inches(6.95),
     "Está apagado por decisión de coordinación (12-ago). Prenderlo cuando se decida toma un minuto: todo el motor ya está construido y probado. Regla de WhatsApp: solo se puede escribir libre si el cliente escribió hace <24 h.", "💡", 1.05)
pie(s, 5)

# ═════════ 6 · SEGUIMIENTO ═════════
s = prs.slides.add_slide(BLANK); cab(s, "04", "Seguimiento de solicitudes", "…/monitor/seguimiento.php")
items = [
    ("🗂️", "Filtros", "Hoy · Ayer · 7 días · 30 días · Histórico — por marca (Ambas / CyR / Carpincentro) y por asesor."),
    ("📅", "Dos fechas, dos preguntas", "«Fecha de entrada» = cuándo escribió el cliente. «Fecha de reporte» = cuándo reportó el asesor. Para «¿qué reportaron ayer?», usa REPORTE."),
    ("📄", "10 por página", "La tabla muestra 10 solicitudes por página (puedes cambiar a 25, 50 o 100) con flechas ‹ › para pasar."),
    ("⬇️", "Exportar a Excel", "Baja exactamente lo que el filtro esté mostrando."),
    ("🎯", "Estados", "🏆 Ganado · 📄 Cotización · 🔄 En gestión · Cerrado · Remitido · ❌ Perdido · ⏳ Pendiente (sin reporte)."),
]
y = 1.85
for e, t_, d_ in items:
    rect(s, Inches(0.75), Inches(y), Inches(11.8), Inches(0.88), fill=CARD, line=LINE, radio=0.12)
    tx(s, Inches(1.0), Inches(y + 0.17), Inches(0.6), Inches(0.5), [[(e, {'s': 20})]])
    tx(s, Inches(1.72), Inches(y + 0.09), Inches(3.4), Inches(0.66), [[(t_, {'s': 14.5, 'b': True, 'c': NAVY})]], anchor=MSO_ANCHOR.MIDDLE)
    tx(s, Inches(5.25), Inches(y + 0.09), Inches(7.1), Inches(0.66), [[(d_, {'s': 12, 'c': SOFT})]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.0
pie(s, 6)

# ═════════ 7 · PANEL WHATSAPP ═════════
s = prs.slides.add_slide(BLANK); cab(s, "05", "Tu panel por WhatsApp", "el 573205662947 no es un cliente para el bot")
cmds = [("informe", "El resumen del día: solicitudes, reparto por asesor, pendientes por reportar y las alertas de los últimos 7 días."),
        ("demo", "Modo clienta: vives el flujo completo como un cliente real, sin tocar datos reales. Ideal para presentaciones."),
        ("hola", "Vuelve al panel (sales del modo demo).")]
y = 2.1
for c_, d_ in cmds:
    rect(s, Inches(1.1), Inches(y), Inches(2.5), Inches(1.0), fill=NAVY, radio=0.14)
    tx(s, Inches(1.1), Inches(y + 0.26), Inches(2.5), Inches(0.5), [[(c_, {'s': 20, 'b': True, 'c': GOLD})]], align=PP_ALIGN.CENTER)
    rect(s, Inches(3.85), Inches(y), Inches(8.4), Inches(1.0), fill=CARD, line=LINE, radio=0.14)
    tx(s, Inches(4.15), Inches(y + 0.14), Inches(7.9), Inches(0.75), [[(d_, {'s': 13.5, 'c': INK})]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.3
pie(s, 7)

# ═════════ 8 · ASESORES ═════════
s = prs.slides.add_slide(BLANK); cab(s, "06", "Cómo reportan los asesores", "dos toques, cero Excel")
pasos = [("1", "La tarjeta", "Le llega el cliente completo: nombre, ciudad, perfil, pedido, fotos y el enlace directo al chat, con el botón «📊 Reportar resultado»."),
         ("2", "El reporte", "Toca el botón y elige: Ganado (con valor), Cotización enviada, En gestión, Cerrado o Perdido."),
         ("3", "El recordatorio", "Si no reporta, el bot le recuerda una vez al día hábil (agrupado, hasta 5 días).")]
x = 0.75
for n_, t_, d_ in pasos:
    rect(s, Inches(x), Inches(1.85), Inches(3.85), Inches(2.15), fill=CARD, line=LINE, radio=0.1)
    badge(s, Inches(x + 0.25), Inches(2.08), n_)
    tx(s, Inches(x + 0.85), Inches(2.1), Inches(2.8), Inches(0.4), [[(t_, {'s': 15.5, 'b': True, 'c': NAVY})]])
    tx(s, Inches(x + 0.28), Inches(2.62), Inches(3.3), Inches(1.3), [[(d_, {'s': 12, 'c': SOFT})]])
    x += 4.05
nota(s, Inches(0.75), Inches(4.35), Inches(11.85),
     "LA REGLA DE ORO — un cliente con solicitud sin reportar que vuelve a escribir SIEMPRE regresa al mismo asesor. Lo garantiza la base de datos, no la memoria del bot.", "📌", 0.75)
nota(s, Inches(0.75), Inches(5.3), Inches(11.85),
     "Ventana de 24 h del asesor: para recibir tarjetas y fotos gratis, debe escribirle al bot de vez en cuando. Si algo se le queda esperando más de 6 horas, el bot LE AVISA SOLO con una plantilla — tocar el botón lo destraba (nuevo, 12-ago).", "📲", 0.95)
pie(s, 8)

# ═════════ 9 · ALERTAS ═════════
s = prs.slides.add_slide(BLANK); cab(s, "07", "Alertas del vigilante", "corre cada hora · graves por correo · todas en tu «informe»")
tabla(s, Inches(0.55), Inches(1.75), Inches(12.25), [
    ["Alerta", "Qué significa", "Qué haces tú"],
    ["cliente_perdido", "Alguien escribió con intención real y no quedó registrado.", "Ábrelo en el panel; si es venta, regístralo o pásaselo a un asesor."],
    ["carrera_consent", "El muro de autorización se repitió tras un «Sí, autorizo». Sev. 3 = grave; sev. 2 = el freno funcionó y el cliente quedó registrado.", "Solo la severidad 3 exige revisar el chat."],
    ["reporte_perdido", "Un asesor recibió «¡Registrado!» pero su reporte no quedó en la base.", "Confírmalo en seguimiento y pide reenviar el reporte."],
    ["cola_adjuntos", "Fotos/audios llevan más de 6 h esperando a un asesor con la ventana cerrada.", "El bot ya le avisó solo al asesor; si persiste al otro día, recuérdaselo tú."],
    ["sin_reportar", "Un asesor acumula 5+ solicitudes sin reporte de más de 10 días.", "Conversación de gestión con ese asesor."],
    ["cron_caido · token_n8n · sesiones_inertes", "Una pieza técnica dejó de funcionar o está por vencer.", "Avisa a soporte técnico."],
], [2.5, 5.55, 4.2], alto=0.78, size=11.5)
pie(s, 9)

# ═════════ 10 · CORREOS ═════════
s = prs.slides.add_slide(BLANK); cab(s, "08", "Correos automáticos", "salen desde noreply@ardisa.com")
tabla(s, Inches(0.75), Inches(1.85), Inches(11.85), [
    ["Correo", "Cuándo", "A quién"],
    ["Reporte semanal Ardisa (Excel)", "Lunes 7:00 a.m.", "Nancy, Paola, María Camila (+ copia oculta a Deicy)"],
    ["Reporte semanal Carpincentro (Excel)", "Lunes 7:00 a.m.", "Paola, María Camila (+ copia oculta a Deicy)"],
    ["Alerta de leads duplicados", "Cuando aparecen", "Solo Deicy"],
    ["Alertas graves del vigilante", "Cuando aparecen", "Solo Deicy"],
], [4.55, 2.5, 4.8], alto=0.66, size=12.5)
nota(s, Inches(0.75), Inches(5.5), Inches(11.85),
     "Si un lunes no llega el Excel: revisa spam primero; si no está, es alerta de cron_caido — avisa a soporte.", "📬", 0.7)
pie(s, 10)

# ═════════ 11 · FASE 2 ═════════
s = prs.slides.add_slide(BLANK); cab(s, "09", "Fase 2: la IA y SAP")
cards = [
    ("✅  Ya en producción", "La IA entiende texto libre y fotos: identifica el producto y la línea, y solo pregunta lo que falta. Los filtros también son inteligentes: proveedores (español e inglés), reclamos → Servicio al Cliente, empleo → canal correcto.", EMERD),
    ("🧪  Cotización con SAP (piloto)", "El bot podrá consultar inventario y disponibilidad directo en SAP. Arranca SIN precios (decisión del 11-ago) y solo para números de prueba, con el interruptor usar_cotiza — se prende sin tocar el bot.", GOLDD),
    ("⏳  Falta solo una cosa", "El token de acceso del servidor SAP (lo entrega el equipo del MCP). Con eso se prende el piloto.", NAVY),
]
y = 1.9
for t_, d_, col in cards:
    rect(s, Inches(0.85), Inches(y), Inches(11.6), Inches(1.45), fill=CARD, line=LINE, radio=0.1)
    rect(s, Inches(0.85), Inches(y), Inches(0.12), Inches(1.45), fill=col, radio=0.5)
    tx(s, Inches(1.25), Inches(y + 0.13), Inches(10.9), Inches(0.42), [[(t_, {'s': 15.5, 'b': True, 'c': col})]])
    tx(s, Inches(1.25), Inches(y + 0.58), Inches(10.9), Inches(0.8), [[(d_, {'s': 12.5, 'c': INK})]])
    y += 1.66
pie(s, 11)

# ═════════ 12 · SI ALGO FALLA ═════════
s = prs.slides.add_slide(BLANK); cab(s, "10", "Si algo falla", "primero lo simple; casi nunca es grave")
tabla(s, Inches(0.55), Inches(1.75), Inches(12.25), [
    ["Síntoma", "Primero revisa", "Si sigue"],
    ["El bot no responde a nadie", "¿El workflow está activo? n8n.ardisa.com → workflow del bot → interruptor «Active».", "Avisa a soporte. Meta reintenta los mensajes: al volver, el bot los atiende."],
    ["El bot responde raro a UN cliente", "Ábrelo en el panel y revisa la conversación completa.", "Reporta el caso con el pantallazo, como siempre."],
    ["El panel no abre", "¿Estás en la red de la oficina o la VPN?", "Soporte técnico."],
    ["Un asesor no recibe tarjetas", "El bot ya le avisa solo; que responda cualquier cosa al bot (ventana 24 h).", "Revisa la alerta cola_adjuntos."],
    ["No llegó el Excel del lunes", "Carpeta de spam.", "Alerta cron_caido + soporte."],
], [3.0, 4.9, 4.35], alto=0.72, size=12)
nota(s, Inches(0.55), Inches(6.2), Inches(12.25),
     "Respaldo: todas las madrugadas (2:30 a.m.) se respalda la base completa y el bot; se guardan 14 días. Nada se pierde aunque el servidor falle.", "🛡️", 0.7)
pie(s, 12)

# ═════════ 13 · CIERRE ═════════
s = prs.slides.add_slide(BLANK); fondo(s, NAVY)
rect(s, Inches(9.3), Inches(4.4), Inches(3.3), Inches(2.3), fill=NAVY2, radio=0.2)
rect(s, Inches(12.2), Inches(6.35), Inches(0.4), Inches(0.4), fill=EMER, forma=MSO_SHAPE.OVAL)
tx(s, Inches(9.3), Inches(5.05), Inches(3.3), Inches(1.2), [[("🤝", {'s': 60, 'c': WHITE})]], align=PP_ALIGN.CENTER)
rect(s, Inches(0.78), Inches(2.35), Inches(2.0), Inches(0.05), fill=GOLD, radio=0.5)
tx(s, Inches(0.75), Inches(2.6), Inches(11.5), Inches(1.8),
   [[("Un asesor que nunca duerme —", {'s': 36, 'b': True, 'c': WHITE})],
    [("y un equipo que nunca pierde un cliente.", {'s': 36, 'b': True, 'c': GOLD})]], sp=2)
tx(s, Inches(0.75), Inches(4.75), Inches(8.2), Inches(1.0),
   [[("Manual de usuario v2 · Grupo Ardisa · agosto de 2026", {'s': 13, 'c': RGBColor(0xC9, 0xDA, 0xD5)})],
    [("El manual técnico (desarrolladores) es aparte: docs/RUNBOOK.md en el repositorio.", {'s': 11, 'c': RGBColor(0x8F, 0xA8, 0xA2)})]])

out = "/tmp/claude-1000/-home-ubuntu-whatsapp-ardisa/aa4cb282-80d9-4597-addc-8c8d448b7df8/scratchpad/Manual-Bot-WhatsApp-Ardisa-v2.pptx"
prs.save(out)
print("OK:", out)
